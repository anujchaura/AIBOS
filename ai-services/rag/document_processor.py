"""
RAG Pipeline – Document Processor
Handles: PDF, DOCX, PPTX, Excel, CSV, Images (OCR), Web URLs, Emails
Uses lazy imports so the service starts even if some optional packages are unavailable.
"""
import os
import asyncio
from pathlib import Path
from typing import Optional

from .chunker import chunk_text
from .embedder import generate_embeddings
from .vector_store import store_chunks


class DocumentProcessor:
    """Processes various document types into chunks ready for RAG"""

    async def process(
        self,
        document_id: str,
        file_path: str,
        file_type: str,
        org_id: str,
        metadata: dict = None,
    ) -> dict:
        """Main entry point – returns processing result"""
        metadata = metadata or {}

        try:
            # Step 1: Extract text
            text, extra_meta = await self._extract_text(file_path, file_type)
            if not text or not text.strip():
                return {"success": False, "error": "No text could be extracted from document"}

            # Step 2: Chunk
            chunks = chunk_text(text, metadata={
                "document_id": document_id,
                "org_id": org_id,
                "file_type": file_type,
                "source": metadata.get("name", file_path),
                **extra_meta,
                **metadata,
            })

            # Step 3: Embed & store
            vector_ids = await store_chunks(chunks, org_id)

            return {
                "success": True,
                "chunk_count": len(chunks),
                "vector_ids": vector_ids,
                "word_count": len(text.split()),
                "char_count": len(text),
                "extracted_text_preview": text[:500],
                **extra_meta,
            }

        except Exception as e:
            return {"success": False, "error": str(e)}

    async def process_url(self, document_id: str, url: str, org_id: str) -> dict:
        """Process a web URL"""
        try:
            import requests
            from bs4 import BeautifulSoup

            headers = {"User-Agent": "Mozilla/5.0 (compatible; AIBOS/1.0)"}
            response = requests.get(url, headers=headers, timeout=30)
            response.raise_for_status()

            soup = BeautifulSoup(response.text, "html.parser")
            for tag in soup(["script", "style", "nav", "footer", "header", "aside"]):
                tag.decompose()

            title = soup.find("title")
            title_text = title.get_text(strip=True) if title else url
            text = soup.get_text(separator="\n", strip=True)

            chunks = chunk_text(text, metadata={
                "document_id": document_id,
                "org_id": org_id,
                "file_type": "url",
                "source": url,
                "title": title_text,
            })

            vector_ids = await store_chunks(chunks, org_id)

            return {
                "success": True,
                "chunk_count": len(chunks),
                "vector_ids": vector_ids,
                "word_count": len(text.split()),
                "title": title_text,
            }
        except Exception as e:
            return {"success": False, "error": str(e)}

    async def _extract_text(self, file_path: str, file_type: str) -> tuple:
        """Extract text based on file type"""
        if file_type == "pdf":
            return await asyncio.get_event_loop().run_in_executor(None, self._extract_pdf, file_path)
        elif file_type == "docx":
            return await asyncio.get_event_loop().run_in_executor(None, self._extract_docx, file_path)
        elif file_type == "pptx":
            return await asyncio.get_event_loop().run_in_executor(None, self._extract_pptx, file_path)
        elif file_type == "xlsx":
            return await asyncio.get_event_loop().run_in_executor(None, self._extract_excel, file_path)
        elif file_type == "csv":
            return await asyncio.get_event_loop().run_in_executor(None, self._extract_csv, file_path)
        elif file_type in ["image", "jpg", "jpeg", "png", "webp"]:
            return await asyncio.get_event_loop().run_in_executor(None, self._extract_image_ocr, file_path)
        elif file_type == "text":
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read(), {}
        else:
            raise ValueError(f"Unsupported file type: {file_type}")

    def _extract_pdf(self, file_path: str) -> tuple:
        try:
            import pymupdf
            doc = pymupdf.open(file_path)
            text_parts = [f"[Page {i + 1}]\n{page.get_text()}" for i, page in enumerate(doc)]
            meta = {"page_count": len(doc)}
            doc.close()
            return "\n\n".join(text_parts), meta
        except ImportError:
            # Fallback: try pypdf2
            try:
                import PyPDF2
                with open(file_path, "rb") as f:
                    reader = PyPDF2.PdfReader(f)
                    text = "\n".join(page.extract_text() or "" for page in reader.pages)
                return text, {"page_count": len(reader.pages)}
            except Exception as e:
                return f"[PDF extraction failed: {e}]", {}

    def _extract_docx(self, file_path: str) -> tuple:
        try:
            from docx import Document as DocxDocument
            doc = DocxDocument(file_path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip():
                        paragraphs.append(row_text)
            return "\n\n".join(paragraphs), {"page_count": len(doc.paragraphs)}
        except Exception as e:
            return f"[DOCX extraction failed: {e}]", {}

    def _extract_pptx(self, file_path: str) -> tuple:
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            slides = []
            for i, slide in enumerate(prs.slides):
                slide_text = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
                slides.append(f"[Slide {i + 1}]\n" + "\n".join(slide_text))
            return "\n\n".join(slides), {"slide_count": len(prs.slides)}
        except Exception as e:
            return f"[PPTX extraction failed: {e}]", {}

    def _extract_excel(self, file_path: str) -> tuple:
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            all_text = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                all_text.append(f"[Sheet: {sheet_name}]")
                for row in ws.iter_rows(values_only=True):
                    row_text = " | ".join(str(c) for c in row if c is not None)
                    if row_text.strip():
                        all_text.append(row_text)
            return "\n\n".join(all_text), {"sheet_count": len(wb.sheetnames)}
        except Exception as e:
            return f"[Excel extraction failed: {e}]", {}

    def _extract_csv(self, file_path: str) -> tuple:
        try:
            import csv
            rows = []
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                reader = csv.reader(f)
                for row in reader:
                    rows.append(" | ".join(row))
            return "\n".join(rows), {"row_count": len(rows)}
        except Exception as e:
            return f"[CSV extraction failed: {e}]", {}

    def _extract_image_ocr(self, file_path: str) -> tuple:
        try:
            import easyocr
            reader = easyocr.Reader(["en"], gpu=False)
            results = reader.readtext(file_path, detail=0)
            text = "\n".join(results)
            return text, {"ocr_engine": "easyocr"}
        except Exception as e:
            return f"[OCR extraction failed: {e}]", {"ocr_engine": "failed"}
